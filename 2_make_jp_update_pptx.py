import os  
import sys  
import time
from openai import AzureOpenAI  
from dotenv import load_dotenv, find_dotenv  
_ = load_dotenv(find_dotenv())  
  
from pydantic import BaseModel  
from openai_utils import get_parsed_completion  
  
from pptx import Presentation  
from pptx.util import Inches, Pt  
from pptx.dml.color import RGBColor  
from pptx.enum.text import MSO_AUTO_SIZE  # 追加  
from pptx.enum.text import PP_ALIGN
  
class UpdateInformation(BaseModel):  
    title: str  
    target: str  
    date: str  
    description: str  
    reference_links: list[str]  
  
def callGPT(system_prompt, user_prompt):  
    messages =  [  
        {'role':'system', 'content':system_prompt},  
        {'role':'user', 'content':user_prompt}  
    ]  
    return get_parsed_completion(messages, UpdateInformation)

def callGPTmock(system_prompt, user_prompt):
    return UpdateInformation(
        title="タイトルほげげげ",
        target="更新対象機能あらららら",
        date="2022-01-01",
        description="更新内容だよーん"*10,
        reference_links=["https://example.com","https://google.com"]
    ),100,100
  
def get_latest_update_file():  
    files = [f for f in os.listdir() if f.startswith("azure_update") and f.endswith(".md")]  
    if not files:  
        raise FileNotFoundError("No azure_update*.md files found.")  
    latest_file = max(files, key=os.path.getctime)  
    return latest_file  

def print_progress(current, total, start_time):
    if current == 1:
        print(f"Processing slide {current}/{total}...")
    else:
        elapsed = time.time() - start_time
        avg = elapsed / (current - 1)
        eta = avg * (total - current + 1)
        print(f"Processing slide {current}/{total}... (avg: {avg:.1f}s/slide, ETA: {eta:.1f}s)")
  
def main():  
    if len(sys.argv) > 1:  
        input_file = sys.argv[1]  
        if not os.path.exists(input_file):  
            print(f"Error: Specified file '{input_file}' not found.")  
            return  
        print(f"Using specified file: {input_file}")  
    else:  
        try:  
            input_file = get_latest_update_file()  
            print(f"Using latest file: {input_file}")  
        except FileNotFoundError as e:  
            print(e)  
            return  
  
    # ファイルを読み込む  
    with open(input_file, "r", encoding="utf-8") as f:  
        update_text = f.read()  
  
    # update_textは =*50で区切られているため分割してループする  
    slides = update_text.split("="*50)  
    
    # Filter out empty slides and count total
    slides = [s for s in slides if s.strip()]
    total_slides = len(slides)
    print(f"Processing {total_slides} slides...")

    system_prompt = """\
入力された情報はあるAzureの機能更新情報です。この情報を以下のルールで更新して指定されたキーのJSONを回答してください  
  
# ルール  
- title: タイトル(先頭行): なるべくそのままの表現で日本語に翻訳する  
- target: 更新対象機能 : 元の情報のまま  
- date: 更新日付: 元の情報のまま  
- description: 更新内容: 日本語で３行程度に要約する  
- reference_links: 参考リンク: 元の情報のまま  
"""  
  
    # プレゼンテーションを作成  
    prs = Presentation()  
  
    # スライドサイズを16:9に設定  
    prs.slide_width = Inches(13.333333)  # 16:9の幅（13.333インチ）  
    prs.slide_height = Inches(7.5)       # 16:9の高さ（7.5インチ）  
  
    start_time = time.time()
    for i, slide_text in enumerate(slides, 1):  
        user_prompt = slide_text.strip()
        print_progress(i, total_slides, start_time)
        # OpenAI APIを呼び出して、情報を取得  
        event, input_token, output_token = callGPT(system_prompt, user_prompt)  
  
        # スライド作成  
        slide_layout = prs.slide_layouts[6]  # 空白レイアウト  
        slide = prs.slides.add_slide(slide_layout)  
  
        # 共通の設定  
        left = Inches(0.5)  
        width = prs.slide_width - Inches(1)  # 左右に0.5インチのマージン  
  
        # 固定の位置を使用  
        top_title = Inches(0.5)  
        top_target = Inches(1.5)  
        top_date = Inches(2.0)  
        top_description = Inches(2.5)  
        top_links = Inches(4.5)  
  
        # タイトルの設定  
        title_box = slide.shapes.add_textbox(left, top_title, width, Inches(1.3))  
        title_frame = title_box.text_frame  
        title_frame.word_wrap = True  
        p = title_frame.paragraphs[0]  
        p.text = event.title  
        p.font.size = Pt(28)  
        p.font.bold = True  
        p.font.name = 'Meiryo UI'  # フォント設定  
        p.font.color.rgb = RGBColor(0, 102, 204)  # 青色

        #タイトルの高さを取得
        title_height_inches = title_box.height / 914400

        # 更新対象機能
        target_top = 0.5 + title_height_inches
        target_box = slide.shapes.add_textbox(left, Inches(target_top), width, Inches(0.5))  
        target_frame = target_box.text_frame  
        target_frame.word_wrap = True  
        p = target_frame.paragraphs[0]  
        p.text = f"対象機能: {event.target}"  
        p.font.size = Pt(16)  
        p.font.name = 'Meiryo UI'  
        p.font.color.rgb = RGBColor(0, 0, 0)  

        # 更新対象機能の高さを取得
        target_height_inches = target_box.height / 914400
  
        # 更新日付
        date_top = target_top + target_height_inches
        date_box = slide.shapes.add_textbox(left, Inches(date_top), width, Inches(0.5))  
        date_frame = date_box.text_frame  
        date_frame.word_wrap = True  
        p = date_frame.paragraphs[0]  
        p.text = f"更新日付: {event.date}"  
        p.font.size = Pt(16)  
        p.font.name = 'Meiryo UI'  
        p.font.color.rgb = RGBColor(0, 0, 0)  

        # 更新日付の高さを取得
        date_height_inches = date_box.height / 914400

        # 更新内容  
        description_top = date_top + date_height_inches
        description_box = slide.shapes.add_textbox(left, Inches(description_top), width, Inches(2.7))  
        description_frame = description_box.text_frame  
        description_frame.word_wrap = True  
        p = description_frame.paragraphs[0]    
        p.text = "内容:"  
        p.font.size = Pt(18)  
        p.font.bold = True  
        p.font.name = 'Meiryo UI'  
        p.font.color.rgb = RGBColor(0, 0, 0)  
  
        # 更新内容の詳細を追加  
        for line in event.description.split("。"):  
            if line.strip():  
                p = description_frame.add_paragraph()  
                p.text = line.strip() + "。"  
                p.font.size = Pt(20)  
                p.font.name = 'Meiryo UI'  
                p.font.color.rgb = RGBColor(0, 0, 0)  

        # 更新内容の高さを取得
        description_height_inches = description_box.height / 914400
  
        # 参考リンク
        link_top = description_top + description_height_inches
        link_box = slide.shapes.add_textbox(left, Inches(link_top), width, Inches(1.5))  
        link_frame = link_box.text_frame  
        link_frame.word_wrap = True  
        p = link_frame.add_paragraph()  
        p.text = "参考リンク:"  
        p.font.size = Pt(18)  
        p.font.bold = True  
        p.font.name = 'Meiryo UI'  
        p.font.color.rgb = RGBColor(0, 0, 0)  
  
        # 各リンクを追加  
        for link in event.reference_links:  
            p = link_frame.add_paragraph()  
            # ハイパーリンクを設定  
            r = p.add_run()
            r.font.size = Pt(14)
            p.font.name = 'Meiryo UI'
            r.text = link
            r.font.color.rgb = RGBColor(0, 0, 255)
            r.hyperlink.address = link  
  
    # PPTXファイルを保存  
    output_file = f"pptx_{os.path.splitext(os.path.basename(input_file))[0]}.pptx"  
    prs.save(output_file)  
    print(f"PPTXファイルを保存しました: {output_file}")  
  
if __name__ == "__main__":  
    main()  